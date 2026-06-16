"""
make_ppt.py — USD/KRW 환율 예측 시스템 기말 보고서 PPT 생성
python3 make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ───────────────────────────────────────────
C_BG       = RGBColor(0x0B, 0x15, 0x28)   # 다크 네이비
C_ACCENT   = RGBColor(0x5E, 0xAE, 0xFF)   # 밝은 파랑
C_GREEN    = RGBColor(0x34, 0xD3, 0x99)   # 에메랄드
C_GOLD     = RGBColor(0xF5, 0x9E, 0x0B)   # 골드
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT    = RGBColor(0xB0, 0xC8, 0xE8)   # 연한 파랑
C_PANEL    = RGBColor(0x0F, 0x2A, 0x4A)   # 패널 배경
C_RED      = RGBColor(0xF8, 0x71, 0x71)
C_MUTED    = RGBColor(0x4A, 0x6A, 0x8A)

W = Inches(13.33)   # 와이드 16:9 너비
H = Inches(7.5)     # 와이드 16:9 높이

# ── 유틸 함수 ─────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # 완전 빈 레이아웃
    return prs.slides.add_slide(layout)

def rect(slide, l, t, w, h, fill=C_BG, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def txt(slide, text, l, t, w, h,
        size=20, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def hline(slide, t, color=C_ACCENT, thickness=Pt(1.5)):
    line = slide.shapes.add_shape(1, Inches(0.5), t, W - Inches(1.0), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def bg(slide):
    rect(slide, 0, 0, W, H, C_BG)

def slide_header(slide, title, subtitle=""):
    bg(slide)
    rect(slide, 0, 0, W, Inches(1.15), C_PANEL)
    txt(slide, title, Inches(0.55), Inches(0.22), Inches(11), Inches(0.7),
        size=32, bold=True, color=C_ACCENT)
    if subtitle:
        txt(slide, subtitle, Inches(0.55), Inches(0.82), Inches(11), Inches(0.35),
            size=14, color=C_LIGHT)
    hline(slide, Inches(1.15))

def card(slide, l, t, w, h, fill=C_PANEL):
    r = rect(slide, l, t, w, h, fill)
    return r

def bullet_box(slide, items, l, t, w, h,
               title="", title_color=C_ACCENT,
               item_color=C_WHITE, size=16, gap=Inches(0.38)):
    if title:
        txt(slide, title, l + Inches(0.12), t + Inches(0.1),
            w - Inches(0.2), Inches(0.4),
            size=15, bold=True, color=title_color)
        t += Inches(0.45)
    for item in items:
        txt(slide, f"▸  {item}", l + Inches(0.15), t, w - Inches(0.3), gap,
            size=size, color=item_color)
        t += gap

def kpi(slide, l, t, w, h, label, value, unit="", val_color=C_GREEN):
    card(slide, l, t, w, h)
    txt(slide, label, l, t + Inches(0.12), w, Inches(0.3),
        size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(slide, value, l, t + Inches(0.38), w, Inches(0.55),
        size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    if unit:
        txt(slide, unit, l, t + Inches(0.88), w, Inches(0.25),
            size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# 슬라이드 제작
# ════════════════════════════════════════════════════════════

prs = new_prs()

# ── Slide 1: 표지 ─────────────────────────────────────────
sl = blank_slide(prs)
bg(sl)
rect(sl, 0, 0, W, Inches(0.08), C_ACCENT)           # 상단 강조선
rect(sl, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT)  # 하단 강조선

txt(sl, "USD/KRW", Inches(1.0), Inches(1.3), Inches(11.3), Inches(1.4),
    size=72, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
txt(sl, "환율 예측 시스템",
    Inches(1.0), Inches(2.55), Inches(11.3), Inches(1.2),
    size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

hline(sl, Inches(3.85), C_MUTED, Pt(1))

txt(sl,
    "ARIMAX · LightGBM · LSTM · BiGRU + Ridge 앙상블\n"
    "멀티스케일 피처(137개) · AWS 자동화 파이프라인 · Streamlit 실시간 대시보드",
    Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.1),
    size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)

txt(sl, "데이터마이닝  ·  기말 보고서",
    Inches(1.0), Inches(5.3), Inches(11.3), Inches(0.4),
    size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
txt(sl, "이용재 (2025720228)  ·  성균관대학교 정보통신대학원 빅데이터학과",
    Inches(1.0), Inches(5.7), Inches(11.3), Inches(0.4),
    size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)
txt(sl, "2026년 1학기",
    Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.4),
    size=13, color=C_MUTED, align=PP_ALIGN.CENTER)

# ── Slide 2: 연구 배경 & 목적 ────────────────────────────
sl = blank_slide(prs)
slide_header(sl, "연구 배경 및 목적",
             "Why USD/KRW?  |  What We Built")
bg_rect_left = card(sl, Inches(0.4), Inches(1.35), Inches(5.9), Inches(5.7))
bg_rect_right = card(sl, Inches(6.6), Inches(1.35), Inches(6.33), Inches(5.7))

txt(sl, "연구 배경", Inches(0.55), Inches(1.45), Inches(5.7), Inches(0.4),
    size=17, bold=True, color=C_ACCENT)
items_bg = [
    "USD/KRW는 수출 의존 경제에서 가장 중요한 거시 변수",
    "기존 연구는 단일 모델·단일 시간 단위에 한정",
    "변동성 레짐 전환 시 예측 오차가 급증하는 구조적 문제",
    "실시간 운영 가능한 엔드투엔드 시스템 부재",
]
bullet_box(sl, items_bg, Inches(0.45), Inches(1.85), Inches(5.8), Inches(5.0),
           item_color=C_LIGHT, size=15, gap=Inches(0.55))

txt(sl, "연구 목적 & 기여", Inches(6.75), Inches(1.45), Inches(6.0), Inches(0.4),
    size=17, bold=True, color=C_GREEN)
items_obj = [
    "D+1 ~ D+22 다중 호라이즌 예측",
    "멀티스케일(일봉·1H·5M·1M) 피처 137개 통합",
    "4개 이종 모델 Ridge 앙상블로 오차 최소화",
    "Lookahead Bias 완전 차단 (시간 역순 분할)",
    "AWS + Streamlit Cloud 실시간 자동화 파이프라인",
    "Persistence Baseline 대비 통계적 유의성 검증",
]
bullet_box(sl, items_obj, Inches(6.65), Inches(1.85), Inches(6.1), Inches(5.0),
           item_color=C_WHITE, size=15, gap=Inches(0.52))

# ── Slide 3: 데이터 수집 & 피처 엔지니어링 ────────────────
sl = blank_slide(prs)
slide_header(sl, "데이터 수집 & 피처 엔지니어링",
             "14 Tickers · 4 시간 단위 · 137 Features")

# 상단 KPI 4개
kw = Inches(2.8)
kh = Inches(1.25)
kpi(sl, Inches(0.4),  Inches(1.35), kw, kh, "학습 기간", "2015–2026", "10년+ 일봉")
kpi(sl, Inches(3.45), Inches(1.35), kw, kh, "수집 티커", "14개", "KRW·DXY·VIX·GOLD 등", C_ACCENT)
kpi(sl, Inches(6.5),  Inches(1.35), kw, kh, "멀티스케일", "4 레이어", "일·1H·5M·1M", C_GOLD)
kpi(sl, Inches(9.55), Inches(1.35), kw, kh, "총 피처 수", "137개", "9개 그룹", C_GREEN)

# 피처 그룹 테이블
headers = ["그룹", "내용", "대표 피처"]
rows = [
    ["G1 추세",   "EMA·SMA·MACD",         "ema_5, macd_cross"],
    ["G2 모멘텀", "RSI·Stoch·CCI·ADX",    "rsi_ob, stoch_sig"],
    ["G3 변동성", "BB·ATR·HV·레짐",       "bb_pos, vol_regime"],
    ["G4 수익률", "다구간 리턴+래그",     "ret_1d~60d, lag_*"],
    ["G5 금리차", "한미 스프레드",         "rate_spread, yield_curve"],
    ["G6 크로스", "CNY·JPY·TWD 상대강도", "krw_vs_cny, cny_corr_20"],
    ["G7 한국",   "SOX·EEM·WTI(KRW)",     "sox_vs_sp, eem_vs_sp"],
    ["G8 멀티스케일", "분봉 통계 정규화", "h1_std_z, ms_vol_ratio"],
    ["G9 캘린더", "요일·월·지정학",       "day_of_week, geo_risk"],
]
col_w = [Inches(1.5), Inches(2.5), Inches(3.0)]
col_x = [Inches(0.4), Inches(1.95), Inches(4.5)]
row_h = Inches(0.37)
top   = Inches(2.75)

# 헤더
for j, (h_lbl, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    r = rect(sl, cx, top, cw - Inches(0.05), row_h, C_ACCENT)
    txt(sl, h_lbl, cx + Inches(0.05), top + Inches(0.05),
        cw - Inches(0.1), row_h - Inches(0.08),
        size=13, bold=True, color=C_BG)

for i, row in enumerate(rows):
    row_top = top + row_h * (i + 1)
    fill = C_PANEL if i % 2 == 0 else RGBColor(0x0D, 0x22, 0x3D)
    for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        rect(sl, cx, row_top, cw - Inches(0.05), row_h - Inches(0.02), fill)
        color = C_GOLD if j == 0 else C_WHITE
        txt(sl, cell, cx + Inches(0.07), row_top + Inches(0.04),
            cw - Inches(0.12), row_h - Inches(0.08),
            size=12, color=color)

# 오른쪽 설명
card(sl, Inches(7.7), Inches(2.75), Inches(5.2), Inches(4.5))
txt(sl, "설계 원칙", Inches(7.85), Inches(2.85), Inches(5.0), Inches(0.4),
    size=15, bold=True, color=C_ACCENT)
principles = [
    "Lookahead Bias 완전 차단",
    "  — 모든 피처는 현재 이전 데이터만 사용",
    "분봉 → 일봉 집계 (std·range·count)",
    "  — 장중 변동성을 일봉에 녹여냄",
    "한국 기준금리 수동 내장 (API 부재)",
    "  — 2015–2026 실제 고시금리 하드코딩",
    "결측치 ffill → bfill → fillna(0)",
    "  — 파이프라인 안정성 확보",
]
t_start = Inches(3.3)
for item in principles:
    color = C_GREEN if "—" not in item else C_LIGHT
    txt(sl, item, Inches(7.88), t_start, Inches(5.0), Inches(0.35),
        size=13, color=color)
    t_start += Inches(0.38)

# ── Slide 4: 모델 아키텍처 ────────────────────────────────
sl = blank_slide(prs)
slide_header(sl, "모델 아키텍처",
             "ARIMAX + LightGBM + LSTM + BiGRU → Ridge 메타 앙상블")

models = [
    ("ARIMAX", C_GOLD,
     ["pmdarima auto_arima", "외생변수 6개", "(rate_spread, DXY, WTI,",
      "VIX, krw_cny, yield)", "→ log_return 예측", "D+1 전용"]),
    ("LightGBM", C_ACCENT,
     ["n_estimators=3,000", "learning_rate=0.02", "num_leaves=63",
      "Early stopping(100)", "→ D+1~D+22 전 호라이즌", "가장 빠른 추론"]),
    ("LSTM", C_GREEN,
     ["30일 슬라이딩 윈도우", "128→64→32 유닛", "Dropout 0.20/0.15/0.10",
      "Huber Loss(δ=0.01)", "Adam(lr=1e-3, clipnorm=1)", "D+1 전용"]),
    ("BiGRU", RGBColor(0xA7, 0x8B, 0xFA),
     ["양방향 GRU 64→32", "Dropout 0.20/0.10",
      "EarlyStopping(15)", "ReduceLROnPlateau",
      "Sharpe 향상 기여", "D+1 전용"]),
]

box_w = Inches(2.85)
box_h = Inches(4.8)
gap   = Inches(0.25)
start = Inches(0.35)

for i, (name, color, items) in enumerate(models):
    lx = start + (box_w + gap) * i
    card(sl, lx, Inches(1.35), box_w, box_h)
    rect(sl, lx, Inches(1.35), box_w, Inches(0.48), color)
    txt(sl, name, lx, Inches(1.42), box_w, Inches(0.4),
        size=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    ty = Inches(1.92)
    for item in items:
        txt(sl, item, lx + Inches(0.12), ty, box_w - Inches(0.2), Inches(0.38),
            size=13, color=C_WHITE)
        ty += Inches(0.46)

# 앙상블 화살표 영역
card(sl, Inches(0.35), Inches(6.3), Inches(11.85), Inches(0.9))
txt(sl, "↓  Ridge 메타 앙상블 (α=0.1)  →  최종 예측 (D+1)",
    Inches(0.35), Inches(6.35), Inches(11.85), Inches(0.55),
    size=17, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
txt(sl, "단순 결합이 아닌 정규화 회귀로 과적합 방지 · 4모델 가중합 최적화",
    Inches(0.35), Inches(6.73), Inches(11.85), Inches(0.4),
    size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

# ── Slide 5: 학습 방법론 ─────────────────────────────────
sl = blank_slide(prs)
slide_header(sl, "학습 방법론",
             "시계열 분할 · Lookahead Bias 차단 · 평가 설계")

# 시계열 분할 시각화
card(sl, Inches(0.4), Inches(1.35), Inches(12.5), Inches(1.5))
txt(sl, "시계열 분할 (70 / 15 / 15)  —  미래 데이터 사용 완전 차단",
    Inches(0.55), Inches(1.42), Inches(12.0), Inches(0.4),
    size=14, bold=True, color=C_ACCENT)
# 분할 바
bar_l, bar_t, bar_w, bar_h = Inches(0.55), Inches(1.88), Inches(12.05), Inches(0.7)
rect(sl, bar_l,                   bar_t, bar_w * 0.70, bar_h, C_ACCENT)
rect(sl, bar_l + bar_w * 0.70,    bar_t, bar_w * 0.15, bar_h, C_GOLD)
rect(sl, bar_l + bar_w * 0.85,    bar_t, bar_w * 0.15, bar_h, C_GREEN)
txt(sl, "Train  70%\n2015.01 – 2023.07",
    bar_l + Inches(0.1), bar_t + Inches(0.08), bar_w * 0.68, bar_h - Inches(0.1),
    size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
txt(sl, "Val  15%\n2023.07 – 2024.06",
    bar_l + bar_w * 0.70 + Inches(0.05), bar_t + Inches(0.08),
    bar_w * 0.13, bar_h - Inches(0.1),
    size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
txt(sl, "Test  15%\n2024.06 – 현재",
    bar_l + bar_w * 0.85 + Inches(0.05), bar_t + Inches(0.08),
    bar_w * 0.13, bar_h - Inches(0.1),
    size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

# 방법론 박스 3개
cols3 = [
    ("평가 지표", C_ACCENT, [
        "RMSE  — 절대 오차 크기 (원화)",
        "MAE   — 평균 절대 오차",
        "MAPE  — 상대 오차율 (%)",
        "DA    — 방향 정확도 (%)",
        "Sharpe — 리스크 조정 수익률",
        "DA_p  — 이항검정 p-value ★",
    ]),
    ("Lookahead Bias 차단", C_GREEN, [
        "Train에만 RobustScaler fit",
        "Val/Test는 transform만 사용",
        "타겟: log(P_{t+h}/P_t) 형태",
        "CLIP_BOUNDS로 이상 예측 억제",
        "피처: 현재 이전 데이터만 사용",
        "미래 정보 0% 보장",
    ]),
    ("한계 & 계획", C_GOLD, [
        "현재: 단순 분할 결과",
        "Walk-Forward Val 미적용",
         "→ 적용 시 성능 변동 가능",
        "VIF 다중공선성 진단 예정",
        "레짐별 앙상블 가설(H3) 검증",
        "기말까지 Walk-Forward 완료 목표",
    ]),
]
bw3 = Inches(4.0)
for i, (title, color, items) in enumerate(cols3):
    lx = Inches(0.4) + (bw3 + Inches(0.25)) * i
    card(sl, lx, Inches(2.95), bw3, Inches(4.3))
    rect(sl, lx, Inches(2.95), bw3, Inches(0.42), color)
    txt(sl, title, lx, Inches(3.0), bw3, Inches(0.38),
        size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    ty = Inches(3.45)
    for item in items:
        c = C_WHITE if "→" not in item else C_GOLD
        txt(sl, item, lx + Inches(0.12), ty, bw3 - Inches(0.2), Inches(0.38),
            size=13, color=c)
        ty += Inches(0.43)

# ── Slide 6: 실험 결과 — 성능 비교 ──────────────────────
sl = blank_slide(prs)
slide_header(sl, "실험 결과 — 성능 비교 (Validation Set)",
             "단순 분할 기준 · Walk-Forward 미적용 · Persistence Baseline 포함")

# 성능 테이블
perf_data = [
    ("Persistence (Baseline)", "8.831", "6.654", "0.505", "49.0", "–", "0.608", C_MUTED),
    ("ARIMAX",                 "8.454", "6.379", "0.484", "47.5", "–", "0.617", C_LIGHT),
    ("LightGBM D+1",           "7.850", "5.874", "0.446", "57.3", "★", "–",     C_ACCENT),
    ("LSTM D+1",               "9.256", "6.962", "0.526", "48.5", "–", "0.348", C_LIGHT),
    ("BiGRU D+1",              "29.55", "27.20", "2.060", "48.8", "–", "0.268", C_RED),
    ("★ Ridge 앙상블 (4모델)", "7.453", "5.812", "0.441", "48.5", "–", "0.367", C_GREEN),
]
col_labels = ["모델", "RMSE(↓)", "MAE(↓)", "MAPE%(↓)", "DA%(↑)", "DA_p<0.05", "Sharpe(↑)"]
col_widths  = [Inches(2.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.15), Inches(1.35), Inches(1.3)]
col_starts  = []
cx = Inches(0.4)
for cw in col_widths:
    col_starts.append(cx)
    cx += cw + Inches(0.04)

rh  = Inches(0.46)
top = Inches(1.38)

# 헤더 행
for label, cs, cw in zip(col_labels, col_starts, col_widths):
    rect(sl, cs, top, cw, rh, C_ACCENT)
    txt(sl, label, cs + Inches(0.05), top + Inches(0.06), cw - Inches(0.08), rh - Inches(0.1),
        size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

for i, (name, rmse, mae, mape, da, sig, sharpe, row_color) in enumerate(perf_data):
    rt = top + rh * (i + 1)
    fill = C_PANEL if i % 2 == 0 else RGBColor(0x0D, 0x22, 0x3D)
    is_best = ("앙상블" in name)
    if is_best:
        fill = RGBColor(0x0A, 0x2A, 0x18)
    for cs, cw in zip(col_starts, col_widths):
        rect(sl, cs, rt, cw, rh - Inches(0.03), fill)
    cells = [name, rmse, mae, mape, da, sig, sharpe]
    for j, (cell, cs, cw) in enumerate(zip(cells, col_starts, col_widths)):
        color = row_color if j == 0 else (C_GOLD if cell == "★" else C_WHITE)
        txt(sl, cell, cs + Inches(0.05), rt + Inches(0.07), cw - Inches(0.08), rh - Inches(0.15),
            size=13, bold=(j == 0 and is_best), color=color, align=PP_ALIGN.CENTER)

# 하단 인사이트
card(sl, Inches(0.4), Inches(5.25), Inches(12.4), Inches(1.95))
txt(sl, "주요 인사이트", Inches(0.55), Inches(5.32), Inches(12.0), Inches(0.38),
    size=14, bold=True, color=C_GOLD)
insights = [
    "LGB RMSE 7.85원으로 단일 모델 중 최고 — 하지만 DA p-value 미검증 시 과대평가 위험",
    "앙상블이 RMSE 7.45원으로 LGB 대비 5% 개선 — Sharpe 0.367로 LSTM(0.348) 상회",
    "BiGRU는 RMSE 29.5원으로 불안정 — 앙상블 내 가중치 낮게 할당됨 (Ridge α=0.1 효과)",
    "DA_p(이항검정) 추가로 '50% 동전던지기 대비 통계적 우위'를 정량 근거로 제시 가능",
]
ty = Inches(5.68)
for ins in insights:
    txt(sl, f"▸  {ins}", Inches(0.55), ty, Inches(12.0), Inches(0.36),
        size=12.5, color=C_LIGHT)
    ty += Inches(0.37)

# ── Slide 7: D+1~D+22 다중 호라이즌 ─────────────────────
sl = blank_slide(prs)
slide_header(sl, "다중 호라이즌 예측 결과 (LightGBM)",
             "D+1 ~ D+22 · Val Set 성능 · 예측 예시 (2026-04-24 기준)")

# 호라이즌 성능 바
horizon_data = [
    (1,   7.85,  57.3, "1478.4",  "상승 ↑"),
    (3,  13.92,  49.8, "1479.4",  "상승 ↑"),
    (5,  16.94,  49.5, "1479.7",  "상승 ↑"),
    (10, 22.45,  47.0, "1478.7",  "상승 ↑"),
    (22, 33.30,  47.5, "1479.3",  "상승 ↑"),
]

card(sl, Inches(0.4), Inches(1.35), Inches(5.5), Inches(5.85))
txt(sl, "호라이즌별 RMSE & 방향 정확도 (Validation)",
    Inches(0.55), Inches(1.42), Inches(5.2), Inches(0.38),
    size=13, bold=True, color=C_ACCENT)

bar_left = Inches(0.55)
max_rmse = 34.0
by = Inches(1.9)
for h, rmse, da, price, direction in horizon_data:
    txt(sl, f"D+{h:2d}",
        bar_left, by, Inches(0.5), Inches(0.35), size=13, color=C_LIGHT)
    bw = Inches(3.0) * (rmse / max_rmse)
    rect(sl, bar_left + Inches(0.55), by + Inches(0.04),
         bw, Inches(0.28), C_ACCENT)
    txt(sl, f"RMSE {rmse:.1f}원 | DA {da:.1f}%",
        bar_left + Inches(0.58), by + Inches(0.06),
        Inches(3.8), Inches(0.26), size=11, color=C_BG)
    by += Inches(0.52)

txt(sl, "※ 호라이즌 증가 → RMSE 급증 / DA는 47~57% 수렴",
    Inches(0.55), Inches(4.0), Inches(5.2), Inches(0.4),
    size=11, color=C_MUTED)

# 오른쪽: 예측 예시
card(sl, Inches(6.2), Inches(1.35), Inches(6.7), Inches(5.85))
txt(sl, "앙상블 예측 예시  (기준: 2026-04-24, 현재가 1477.96원)",
    Inches(6.35), Inches(1.42), Inches(6.4), Inches(0.38),
    size=13, bold=True, color=C_GREEN)

forecast_rows = [
    ("D+1",  "1,478.4원", "+0.03%", "상승 ↑", "LGB/ARIMAX/LSTM/BiGRU"),
    ("D+3",  "1,479.4원", "+0.10%", "상승 ↑", "LGB 단독"),
    ("D+5",  "1,479.7원", "+0.12%", "상승 ↑", "LGB 단독"),
    ("D+10", "1,478.7원", "+0.05%", "상승 ↑", "LGB 단독"),
    ("D+22", "1,479.3원", "+0.09%", "상승 ↑", "LGB 단독"),
]
fy = Inches(1.88)
for horizon, price, chg, direc, note in forecast_rows:
    card(sl, Inches(6.3), fy, Inches(6.5), Inches(0.82))
    txt(sl, horizon, Inches(6.45), fy + Inches(0.08), Inches(0.65), Inches(0.38),
        size=15, bold=True, color=C_ACCENT)
    txt(sl, price,   Inches(7.15), fy + Inches(0.08), Inches(1.8),  Inches(0.38),
        size=15, bold=True, color=C_WHITE)
    txt(sl, chg,     Inches(9.0),  fy + Inches(0.08), Inches(1.0),  Inches(0.38),
        size=14, color=C_GREEN)
    txt(sl, direc,   Inches(10.1), fy + Inches(0.08), Inches(1.2),  Inches(0.38),
        size=13, color=C_GREEN)
    txt(sl, note,    Inches(6.45), fy + Inches(0.46), Inches(5.8),  Inches(0.3),
        size=10, color=C_MUTED)
    fy += Inches(0.93)

# ── Slide 8: Ablation & 통계적 유의성 ───────────────────
sl = blank_slide(prs)
slide_header(sl, "Ablation Study & 통계적 유의성 검증",
             "교수 피드백 반영 — Baseline 비교 + 이항검정 p-value")

# 왼쪽: Ablation
card(sl, Inches(0.4), Inches(1.35), Inches(6.2), Inches(5.9))
txt(sl, "D+1 Ablation Table",
    Inches(0.55), Inches(1.42), Inches(6.0), Inches(0.38),
    size=15, bold=True, color=C_ACCENT)

abl_rows = [
    ("Persistence",     "8.831", "49.0%", "0.608", C_MUTED,   "Baseline"),
    ("ARIMAX",          "8.454", "47.5%", "0.617", C_LIGHT,   "약한 개선"),
    ("LSTM",            "9.256", "48.5%", "0.348", C_LIGHT,   "Sharpe 기여"),
    ("BiGRU",           "29.55", "48.8%", "0.268", C_RED,     "불안정"),
    ("LGB",             "7.850", "57.3%", "—",     C_ACCENT,  "DA 최고"),
    ("★앙상블 (4모델)", "7.453", "48.5%", "0.367", C_GREEN,   "RMSE 최고"),
]
abl_labels = ["모델", "RMSE(↓)", "DA%(↑)", "Sharpe(↑)", "특이사항"]
abl_widths  = [Inches(1.65), Inches(0.95), Inches(0.85), Inches(0.9), Inches(1.5)]
axs = []
ax = Inches(0.45)
for aw in abl_widths:
    axs.append(ax)
    ax += aw + Inches(0.03)

arh = Inches(0.42)
aty = Inches(1.85)
for j, (lbl, ax2, aw) in enumerate(zip(abl_labels, axs, abl_widths)):
    rect(sl, ax2, aty, aw, arh, C_PANEL)
    txt(sl, lbl, ax2 + Inches(0.04), aty + Inches(0.06),
        aw - Inches(0.06), arh - Inches(0.1),
        size=12, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

for i, (name, rmse, da, sharpe, ncolor, note) in enumerate(abl_rows):
    ry = aty + arh * (i + 1) + Inches(0.02)
    fill = RGBColor(0x0A, 0x2A, 0x18) if "앙상블" in name else \
           (C_PANEL if i % 2 == 0 else RGBColor(0x0D, 0x22, 0x3D))
    cells = [name, rmse, da, sharpe, note]
    for j, (cell, ax2, aw) in enumerate(zip(cells, axs, abl_widths)):
        rect(sl, ax2, ry, aw, arh - Inches(0.03), fill)
        color = ncolor if j == 0 else (C_GREEN if j == 4 and "최고" in cell else C_WHITE)
        txt(sl, cell, ax2 + Inches(0.04), ry + Inches(0.06),
            aw - Inches(0.06), arh - Inches(0.1),
            size=12, color=color, align=PP_ALIGN.CENTER)

# 오른쪽: 이항검정 설명
card(sl, Inches(6.85), Inches(1.35), Inches(6.05), Inches(5.9))
txt(sl, "이항검정으로 통계적 유의성 검증",
    Inches(7.0), Inches(1.42), Inches(5.8), Inches(0.38),
    size=15, bold=True, color=C_GREEN)

binom_content = [
    ("검정 설계", C_ACCENT, [
        "H₀ : DA = 50% (동전던지기와 동일)",
        "H₁ : DA > 50% (단측 검정)",
        "검정통계량 : Binomial Test",
        "scipy.stats.binomtest(n_correct, n, p=0.5)",
        "유의수준 : α = 0.05",
    ]),
    ("해석", C_GOLD, [
        "p < 0.05 → 동전던지기 대비 유의한 우위 ★",
        "p ≥ 0.05 → 우연 범위 내, 결론 보류",
        "51~57% DA는 표본오차 내 일 수 있음",
        "재학습 후 DA_p 컬럼 자동 생성됨",
    ]),
    ("향후 계획", C_MUTED, [
        "Walk-Forward Validation 적용",
        "→ 이동 창(expanding window) 검증",
        "레짐별 DA_p 분해 (고변동 vs 저변동)",
        "부트스트랩 신뢰구간 추가 예정",
    ]),
]
ty2 = Inches(1.9)
for section_title, st_color, items in binom_content:
    txt(sl, section_title, Inches(7.0), ty2, Inches(5.8), Inches(0.35),
        size=13, bold=True, color=st_color)
    ty2 += Inches(0.37)
    for item in items:
        txt(sl, f"  {item}", Inches(7.05), ty2, Inches(5.75), Inches(0.33),
            size=12, color=C_WHITE if "→" not in item else C_GOLD)
        ty2 += Inches(0.35)
    ty2 += Inches(0.15)

# ── Slide 9: 배포 아키텍처 ──────────────────────────────
sl = blank_slide(prs)
slide_header(sl, "실시간 배포 아키텍처",
             "Colab GPU → GitHub → AWS → Streamlit Cloud  |  End-to-End 자동화")

components = [
    (Inches(0.4),  Inches(1.4),  Inches(2.5), Inches(2.5), C_GOLD,
     "Google Colab", ["T4 GPU", "train.py", "주 1회 수동", "→ outputs/ 생성"]),
    (Inches(3.15), Inches(1.4),  Inches(2.5), Inches(2.5), C_ACCENT,
     "GitHub", ["모델(.pkl/.keras)", "코드 버전 관리", "Actions: s3-sync", "UTC 21:30 스케줄"]),
    (Inches(5.9),  Inches(1.4),  Inches(2.5), Inches(2.5), C_GREEN,
     "AWS S3", ["forecast_today.json", "residual_data.csv", "모델 파일 백업", "boto3 자동 업로드"]),
    (Inches(8.65), Inches(1.4),  Inches(2.5), Inches(2.5), RGBColor(0xA7, 0x8B, 0xFA),
     "EC2 / Lambda", ["predict.py 일일실행", "KST 06:00 cron", "잔차 수집", "D+1 오차 누적"]),
    (Inches(11.4), Inches(1.4),  Inches(1.5), Inches(2.5), RGBColor(0xFF, 0x4B, 0x4B),
     "Streamlit", ["실시간 차트", "LGB 추론", "성능 대시보드", "5분 캐시"]),
]

for lx, lt, lw, lh, color, name, items in components:
    card(sl, lx, lt, lw, lh)
    rect(sl, lx, lt, lw, Inches(0.45), color)
    txt(sl, name, lx, lt + Inches(0.06), lw, Inches(0.38),
        size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    ty3 = lt + Inches(0.55)
    for item in items:
        txt(sl, item, lx + Inches(0.1), ty3, lw - Inches(0.15), Inches(0.45),
            size=12, color=C_WHITE)
        ty3 += Inches(0.45)

# 화살표 라벨
arrows = [
    (Inches(2.92), Inches(2.55), "git push"),
    (Inches(5.67), Inches(2.55), "S3 sync"),
    (Inches(8.42), Inches(2.55), "cron"),
    (Inches(11.17), Inches(2.55), "webhook"),
]
for alx, aly, label in arrows:
    txt(sl, f"→\n{label}", alx, aly, Inches(0.7), Inches(0.6),
        size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

# Colab 자동화 노트북 셀 구조
card(sl, Inches(0.4), Inches(4.1), Inches(12.5), Inches(3.15))
txt(sl, "colab_master.ipynb — 8셀 자동화 파이프라인",
    Inches(0.55), Inches(4.18), Inches(12.0), Inches(0.38),
    size=14, bold=True, color=C_ACCENT)

cells_info = [
    ("셀 1", "GitHub Secrets 로드", C_GOLD),
    ("셀 2", "Drive 마운트 + 레포 클론", C_LIGHT),
    ("셀 3", "패키지 설치 (requirements_train.txt)", C_LIGHT),
    ("셀 4", "git pull (최신 코드)", C_LIGHT),
    ("셀 5", "train.py or predict.py 자동 판단 실행", C_GREEN),
    ("셀 6", "Streamlit 로컬 테스트 (선택)", C_MUTED),
    ("셀 7", "GitHub push (outputs/ 포함)", C_ACCENT),
    ("셀 8", "스케줄러 (매일 KST 06:00)", C_GOLD),
]
cy3 = Inches(4.6)
cx3 = Inches(0.5)
for i, (cell_no, desc, color) in enumerate(cells_info):
    col = i % 4
    row = i // 4
    lx3 = cx3 + Inches(3.12) * col
    ly3 = cy3 + Inches(1.0) * row
    card(sl, lx3, ly3, Inches(3.0), Inches(0.82))
    txt(sl, cell_no, lx3 + Inches(0.1), ly3 + Inches(0.08), Inches(0.6), Inches(0.38),
        size=13, bold=True, color=color)
    txt(sl, desc, lx3 + Inches(0.7), ly3 + Inches(0.1), Inches(2.2), Inches(0.55),
        size=11, color=C_WHITE)

# ── Slide 10: 결론 & 향후 계획 ──────────────────────────
sl = blank_slide(prs)
slide_header(sl, "결론 및 향후 계획",
             "What We Achieved  ·  What's Next")

card(sl, Inches(0.4), Inches(1.35), Inches(5.85), Inches(5.85))
txt(sl, "달성한 것", Inches(0.55), Inches(1.42), Inches(5.6), Inches(0.38),
    size=16, bold=True, color=C_GREEN)
achieved = [
    ("✅ 4모델 앙상블 구축",         "ARIMAX+LGB+LSTM+BiGRU+Ridge"),
    ("✅ RMSE 7.45원 달성",          "LGB 단독(7.85원) 대비 5% 개선"),
    ("✅ D+1~D+22 다중 호라이즌",    "실서비스 수준의 예측 범위"),
    ("✅ 137개 멀티스케일 피처",     "분봉 통계를 일봉에 통합"),
    ("✅ End-to-End 자동화",         "Colab→GitHub→AWS→Streamlit"),
    ("✅ Lookahead Bias 완전 차단",  "시간 역순 분할 + 스케일러 분리"),
    ("✅ 이항검정 p-value 추가",     "교수 피드백 반영 (DA 유의성)"),
    ("✅ Persistence Baseline",      "기준선 대비 성능 정량 비교"),
]
ty4 = Inches(1.85)
for title, desc in achieved:
    txt(sl, title, Inches(0.55), ty4, Inches(3.5), Inches(0.35),
        size=13, bold=True, color=C_GREEN)
    txt(sl, desc, Inches(4.15), ty4, Inches(1.9), Inches(0.35),
        size=11, color=C_LIGHT)
    ty4 += Inches(0.42)

card(sl, Inches(6.5), Inches(1.35), Inches(6.4), Inches(5.85))
txt(sl, "향후 계획 (기말까지)", Inches(6.65), Inches(1.42), Inches(6.1), Inches(0.38),
    size=16, bold=True, color=C_GOLD)

plans = [
    ("Walk-Forward Validation",
     "Expanding window으로 일반화 성능 검증\nDA_p가 Walk-Forward에서도 유의한지 확인"),
    ("VIF 다중공선성 진단",
     "137개 피처 상관 구조 분석\n|r|>0.8 페어 파악 및 피처 정제"),
    ("레짐별 앙상블 (H3 가설)",
     "고변동/저변동 레짐 분류 후 별도 모델 적용\n레짐별 RMSE 분해 정량 보고"),
    ("잔차 보정 모델",
     "residual_data.csv 누적 후 학습\n24h 예측 오차의 패턴 포착"),
]
ty5 = Inches(1.85)
for i, (ptitle, pdesc) in enumerate(plans):
    card(sl, Inches(6.55), ty5, Inches(6.2), Inches(1.28))
    txt(sl, f"{i+1}. {ptitle}", Inches(6.7), ty5 + Inches(0.08),
        Inches(6.0), Inches(0.38), size=14, bold=True, color=C_GOLD)
    txt(sl, pdesc, Inches(6.7), ty5 + Inches(0.45),
        Inches(6.0), Inches(0.72), size=12, color=C_LIGHT)
    ty5 += Inches(1.38)

# 하단 quote
rect(sl, Inches(0.4), Inches(7.1), W - Inches(0.8), Inches(0.22), C_ACCENT)
txt(sl,
    '"이미 구축한 파이프라인의 결과가 얼마나 신뢰할 만한가를 정량적으로 검증하는 데 무게를 두어야 한다."  — 구자환 교수',
    Inches(0.5), Inches(7.14), W - Inches(1.0), Inches(0.28),
    size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
out = "/home/user/usdkrw-prediction/usdkrw_final_report.pptx"
prs.save(out)
print(f"✅ 저장 완료: {out}")
print(f"   슬라이드 수: {len(prs.slides)}")
