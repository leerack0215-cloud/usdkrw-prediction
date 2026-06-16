"""
make_ppt_v2.py  —  USD/KRW 환율 예측 시스템 (5분 발표, 8슬라이드)
레퍼런스 스타일: 다크 네이비 + 그린 포인트 / 컬러 테두리 카드
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 ──────────────────────────────────────────────────
C_BG    = RGBColor(0x0B, 0x11, 0x20)
C_PANEL = RGBColor(0x0F, 0x1E, 0x30)
C_GREEN = RGBColor(0x22, 0xC5, 0x5E)
C_CYAN  = RGBColor(0x22, 0xD3, 0xEE)
C_GOLD  = RGBColor(0xF5, 0x9E, 0x0B)
C_RED   = RGBColor(0xEF, 0x44, 0x44)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED = RGBColor(0x64, 0x74, 0x8B)
C_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
C_DARK2 = RGBColor(0x13, 0x27, 0x40)

W = Inches(13.33)
H = Inches(7.5)

# ── 유틸 ──────────────────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(sl):
    s = sl.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = C_BG
    s.line.fill.background()

def box(sl, l, t, w, h, fill=C_PANEL, border=None, bw=Pt(1.5)):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def txt(sl, text, l, t, w, h,
        size=16, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def section_label(sl, label):
    txt(sl, label, Inches(0.45), Inches(0.2), Inches(12), Inches(0.35),
        size=9, color=C_GREEN, italic=False)

def slide_title(sl, big, small=""):
    txt(sl, big, Inches(0.45), Inches(0.52), Inches(12), Inches(0.75),
        size=38, bold=True, color=C_WHITE)
    if small:
        txt(sl, small, Inches(0.45), Inches(1.22), Inches(12), Inches(0.3),
            size=13, color=C_LIGHT)

def hline(sl, t, color=C_MUTED):
    s = sl.shapes.add_shape(1, Inches(0.45), t, W - Inches(0.9), Pt(1.5))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()

def card3(sl, lx, t, w, h, border_color, title, title_color, body_lines, body_size=13):
    """3컬럼 카드 (컬러 테두리)"""
    box(sl, lx, t, w, h, fill=C_PANEL, border=border_color, bw=Pt(1.5))
    txt(sl, title, lx + Inches(0.15), t + Inches(0.12), w - Inches(0.25), Inches(0.38),
        size=14, bold=True, color=title_color)
    ty = t + Inches(0.52)
    for line in body_lines:
        txt(sl, line, lx + Inches(0.15), ty, w - Inches(0.25), Inches(0.35),
            size=body_size, color=C_LIGHT)
        ty += Inches(0.35)

def highlight_box(sl, t, h, emoji, bold_text, sub_text, border=C_GREEN):
    """레퍼런스의 전체너비 강조 박스"""
    box(sl, Inches(0.45), t, W - Inches(0.9), h, fill=C_PANEL, border=border, bw=Pt(2))
    txt(sl, f"{emoji}  {bold_text}", Inches(0.65), t + Inches(0.12),
        W - Inches(1.3), Inches(0.38), size=15, bold=True, color=C_WHITE)
    txt(sl, sub_text, Inches(0.65), t + Inches(0.5), W - Inches(1.3), Inches(0.4),
        size=13, color=C_LIGHT)

def bar_row(sl, lx, ty, bw_max, label, value, max_val, bar_color, label_w=Inches(1.2)):
    txt(sl, label, lx, ty, label_w, Inches(0.32), size=12, color=C_LIGHT)
    bar_len = bw_max * (value / max_val)
    b = sl.shapes.add_shape(1, lx + label_w + Inches(0.08), ty + Inches(0.04),
                             bar_len, Inches(0.25))
    b.fill.solid(); b.fill.fore_color.rgb = bar_color; b.line.fill.background()
    txt(sl, f"{value}", lx + label_w + bar_len + Inches(0.12), ty,
        Inches(1.0), Inches(0.32), size=11, color=C_WHITE)

def kpi_box(sl, lx, t, w, h, label, value, unit="", val_color=C_GREEN):
    box(sl, lx, t, w, h, fill=C_PANEL, border=val_color, bw=Pt(1.5))
    txt(sl, label, lx, t + Inches(0.1), w, Inches(0.3),
        size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)
    txt(sl, value, lx, t + Inches(0.38), w, Inches(0.52),
        size=24, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    if unit:
        txt(sl, unit, lx, t + Inches(0.9), w, Inches(0.22),
            size=10, color=C_MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
prs = new_prs()

# ── SL 1: 표지 ───────────────────────────────────────────
sl = blank(prs); bg(sl)

# 왼쪽 세로 그린 바
s = sl.shapes.add_shape(1, 0, 0, Inches(0.12), H)
s.fill.solid(); s.fill.fore_color.rgb = C_GREEN; s.line.fill.background()

txt(sl, "USD/KRW", Inches(0.55), Inches(0.9), Inches(12), Inches(1.5),
    size=80, bold=True, color=C_GREEN)
txt(sl, "환율 예측 시스템", Inches(0.55), Inches(2.2), Inches(12), Inches(1.0),
    size=44, bold=True, color=C_WHITE)
txt(sl, "기말 발표", Inches(0.55), Inches(3.2), Inches(4), Inches(0.42),
    size=18, color=C_LIGHT)
hline(sl, Inches(3.72), C_MUTED)
txt(sl,
    "ARIMAX · LightGBM · LSTM · BiGRU + Ridge 앙상블\n"
    "137개 멀티스케일 피처  ·  AWS 자동화  ·  Streamlit 실시간 대시보드",
    Inches(0.55), Inches(3.85), Inches(12), Inches(0.9),
    size=16, color=C_LIGHT)

# 하단 정보 바
box(sl, 0, Inches(6.85), W, Inches(0.65), fill=C_DARK2)
txt(sl, "정보통신대학원 빅데이터학과  ·  2025720228  ·  이용재",
    Inches(0.5), Inches(6.93), Inches(8), Inches(0.42),
    size=14, color=C_LIGHT)
txt(sl, "github.com/leerack0215-cloud/usdkrw-prediction",
    Inches(0.5), Inches(6.93), W - Inches(0.8), Inches(0.42),
    size=12, color=C_GREEN, align=PP_ALIGN.RIGHT)

# ── SL 2: 연구 개요 ──────────────────────────────────────
sl = blank(prs); bg(sl)
section_label(sl, "R E S E A R C H   O V E R V I E W  ·  연 구 개 요")
slide_title(sl, "연구 개요 & 연구 질문")
hline(sl, Inches(1.52))

highlight_box(sl, Inches(1.62), Inches(0.75),
    "💡", "연구 질문",
    "환율 예측에서 멀티스케일 피처(137개) + 4모델 앙상블이 단일 모델보다 통계적으로 우월한가?\n"
    "그리고 DA(방향 정확도)는 동전던지기(50%) 대비 유의미한가?")

C_BLUE = RGBColor(0x3B, 0x82, 0xF6)
bw = Inches(3.9); gap = Inches(0.27)
bx = Inches(0.45); by = Inches(2.55); bh = Inches(2.05)

card3(sl, bx, by, bw, bh, C_GREEN,  "📌  배경",   C_GREEN,
    ["USD/KRW는 수출 의존 경제의 핵심 거시변수",
     "기존 연구: 단일 모델·단일 시간 단위",
     "변동성 레짐 전환 시 예측 오차 급증"],
    body_size=12)
card3(sl, bx + bw + gap, by, bw, bh, C_CYAN, "💡  목표",   C_CYAN,
    ["D+1 ~ D+22 다중 호라이즌 예측",
     "4개 이종 모델 Ridge 앙상블",
     "Lookahead Bias 완전 차단",
     "통계적 유의성 검증(이항검정)"],
    body_size=12)
card3(sl, bx + (bw + gap)*2, by, bw, bh, C_GOLD, "★  중간→기말",C_GOLD,
    ["중간: LGB 단독 · 단순 분할",
     "기말: 4모델 앙상블 · Ridge 메타러너",
     "     Persistence Baseline 추가",
     "     이항검정 p-value 도입"],
    body_size=12)

txt(sl, "핵심: 단순 분할 결과 기준 — Walk-Forward Validation 적용 시 수치 변동 가능",
    Inches(0.45), Inches(4.75), W - Inches(0.9), Inches(0.3),
    size=11, color=C_GREEN, italic=True)

# ── SL 3: 데이터 & 피처 ──────────────────────────────────
sl = blank(prs); bg(sl)
section_label(sl, "D A T A   C O L L E C T I O N  ·  데 이 터   수 집")
slide_title(sl, "멀티스케일 데이터 & 피처 137개")
hline(sl, Inches(1.52))

# 파이프라인 flow
flow_items = [
    ("yfinance\n14 Tickers", C_GREEN),
    ("일봉+1H\n+5M+1M", C_CYAN),
    ("make_features\n137 피처", C_GOLD),
    ("LGB·LSTM\nBiGRU·ARIMAX", C_GREEN),
    ("Ridge\n앙상블", C_RED),
]
fw = Inches(2.2); fh = Inches(0.88); fy = Inches(1.62)
fx = Inches(0.45)
for i, (label, col) in enumerate(flow_items):
    box(sl, fx, fy, fw, fh, fill=C_PANEL, border=col, bw=Pt(1.5))
    txt(sl, label, fx, fy + Inches(0.08), fw, fh - Inches(0.12),
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        txt(sl, "→", fx + fw, fy + Inches(0.28), Inches(0.32), Inches(0.35),
            size=18, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)
    fx += fw + Inches(0.32)

# KPI 4개
kw = Inches(2.85); ky = Inches(2.62); kh = Inches(1.15)
kpi_box(sl, Inches(0.45), ky, kw, kh, "학습 기간", "10년+", "2015–2026 일봉", C_GREEN)
kpi_box(sl, Inches(0.45)+kw+Inches(0.18), ky, kw, kh, "수집 티커", "14개", "KRW·DXY·VIX 등", C_CYAN)
kpi_box(sl, Inches(0.45)+(kw+Inches(0.18))*2, ky, kw, kh, "시간 단위", "4 레이어", "일·1H·5M·1M", C_GOLD)
kpi_box(sl, Inches(0.45)+(kw+Inches(0.18))*3, ky, kw, kh, "총 피처", "137개", "9개 그룹", C_GREEN)

# 피처 그룹 미니 테이블
groups = [
    ("G1 추세",   "EMA·SMA·MACD",      "16개"),
    ("G2 모멘텀", "RSI·Stoch·CCI",     "11개"),
    ("G3 변동성", "BB·ATR·HV·레짐",    "15개"),
    ("G4 수익률", "다구간 리턴+래그",  "12개"),
    ("G5 금리차", "한미 스프레드",      "8개"),
    ("G6 크로스", "CNY·JPY·TWD",       "9개"),
    ("G7 한국",   "SOX·EEM·WTI",       "12개"),
    ("G8 멀티스케일","분봉 통계 정규화","17개"),
    ("G9 캘린더", "요일·월·지정학",    "6개"),
]
col_w = [Inches(1.5), Inches(2.5), Inches(0.8)]
col_x = [Inches(0.45), Inches(2.0), Inches(4.55)]
rh = Inches(0.32); gy = Inches(3.9)

# 헤더
for lbl, cx, cw in zip(["그룹","내용","수"], col_x, col_w):
    box(sl, cx, gy, cw - Inches(0.04), rh, fill=C_GREEN)
    txt(sl, lbl, cx + Inches(0.05), gy + Inches(0.04), cw - Inches(0.1), rh,
        size=11, bold=True, color=C_BG)
for i, (g, c, n) in enumerate(groups):
    ry = gy + rh * (i + 1) + Inches(0.02)
    fill = C_PANEL if i % 2 == 0 else C_DARK2
    for lbl, cx, cw in zip([g, c, n], col_x, col_w):
        box(sl, cx, ry, cw - Inches(0.04), rh, fill=fill)
        color = C_GREEN if cx == col_x[0] else C_WHITE
        txt(sl, lbl, cx + Inches(0.05), ry + Inches(0.04),
            cw - Inches(0.1), rh, size=10, color=color)

# 우측 설계 원칙
px = Inches(5.6); pw = Inches(7.5)
box(sl, px, Inches(2.55), pw, Inches(4.7), fill=C_PANEL, border=C_MUTED, bw=Pt(1))
txt(sl, "🔒  설계 원칙", px + Inches(0.2), Inches(2.65), pw, Inches(0.38),
    size=14, bold=True, color=C_GREEN)
principles = [
    ("Lookahead Bias 완전 차단",  "Train에만 RobustScaler fit / Val·Test는 transform"),
    ("멀티스케일 → 일봉 통합",    "1H·5M·1M 분봉 통계(std·range·count) 집계"),
    ("결측치 안전 처리",           "ffill → bfill → fillna(0) 3단계"),
    ("이상치 억제",                "log-return CLIP_BOUNDS 적용"),
]
py = Inches(3.1)
for title, desc in principles:
    txt(sl, f"▸  {title}", px + Inches(0.2), py, pw - Inches(0.4), Inches(0.3),
        size=13, bold=True, color=C_CYAN)
    txt(sl, f"     {desc}", px + Inches(0.2), py + Inches(0.32), pw - Inches(0.4), Inches(0.3),
        size=12, color=C_LIGHT)
    py += Inches(0.72)

# ── SL 4: 모델 아키텍처 ──────────────────────────────────
sl = blank(prs); bg(sl)
section_label(sl, "M O D E L   A R C H I T E C T U R E  ·  모 델   구 조")
slide_title(sl, "4모델 앙상블 구조")
hline(sl, Inches(1.52))

models = [
    ("ARIMAX",       C_GOLD,
     ["pmdarima auto_arima", "외생변수 6개",
      "rate_spread · DXY", "WTI · VIX · yield",
      "→ log_return 예측", "D+1 전용"]),
    ("LightGBM",     C_GREEN,
     ["n_estimators 3,000", "learning_rate 0.02",
      "num_leaves 63", "Early stopping(100)",
      "→ D+1 ~ D+22", "가장 빠른 추론"]),
    ("LSTM",         C_CYAN,
     ["30일 슬라이딩 윈도우", "128→64→32 유닛",
      "Dropout 0.20·0.15·0.10", "Huber Loss(δ=0.01)",
      "Adam clipnorm=1", "D+1 전용"]),
    ("BiGRU",        RGBColor(0xA7, 0x8B, 0xFA),
     ["양방향 GRU 64→32", "Dropout 0.20·0.10",
      "EarlyStopping(15)", "ReduceLROnPlateau",
      "Sharpe 향상 기여", "D+1 전용"]),
]

mw = Inches(2.9); mh = Inches(4.6); mx = Inches(0.45)
for name, col, items in models:
    box(sl, mx, Inches(1.62), mw, mh, fill=C_PANEL, border=col, bw=Pt(2))
    box(sl, mx, Inches(1.62), mw, Inches(0.5), fill=col)
    txt(sl, name, mx, Inches(1.67), mw, Inches(0.42),
        size=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    ty = Inches(2.22)
    for item in items:
        txt(sl, item, mx + Inches(0.15), ty, mw - Inches(0.25), Inches(0.38),
            size=12, color=C_WHITE)
        ty += Inches(0.42)
    mx += mw + Inches(0.22)

# 앙상블 박스
box(sl, Inches(0.45), Inches(6.3), W - Inches(0.9), Inches(0.88),
    fill=C_DARK2, border=C_GREEN, bw=Pt(2))
txt(sl, "↓   Ridge 메타 앙상블 (α = 0.1)   →   최종 예측 D+1",
    Inches(0.45), Inches(6.38), W - Inches(0.9), Inches(0.45),
    size=18, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
txt(sl, "4모델 예측값을 입력 → Ridge 정규화 회귀로 최적 가중합 학습 → 과적합 방지",
    Inches(0.45), Inches(6.76), W - Inches(0.9), Inches(0.3),
    size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

# ── SL 5: 실험 결과 ───────────────────────────────────────
sl = blank(prs); bg(sl)
section_label(sl, "R E S U L T S  ·  실 험   결 과")
slide_title(sl, "Ablation Study & 통계적 유의성 검증")
hline(sl, Inches(1.52))

# 성능 테이블 (왼쪽)
perf = [
    ("Persistence",      "8.83", "49.0%", "—",    "0.608", C_MUTED,  "기준선"),
    ("ARIMAX",           "8.45", "47.5%", "—",    "−0.095",C_LIGHT,  "선형 시계열"),
    ("LSTM",             "9.26", "48.5%", "—",    "0.348", C_LIGHT,  "Sharpe 기여"),
    ("BiGRU",            "29.6", "48.8%", "—",    "0.268", C_RED,    "불안정"),
    ("LightGBM",         "7.85", "57.3%", "★",   "−0.102",C_GREEN,  "DA 최고"),
    ("★ 앙상블(4모델)", "7.45", "48.5%", "—",    "0.367", C_GREEN,  "RMSE 최고"),
]
headers = ["모델", "RMSE↓", "DA%↑", "DA_p", "Sharpe↑", "특이사항"]
cws = [Inches(2.0), Inches(0.9), Inches(0.85), Inches(0.75), Inches(0.95), Inches(1.35)]
cxs = []
cx = Inches(0.45)
for cw in cws:
    cxs.append(cx); cx += cw + Inches(0.04)

rh = Inches(0.44); ty = Inches(1.62)
for j, (h, cx2, cw) in enumerate(zip(headers, cxs, cws)):
    box(sl, cx2, ty, cw, rh, fill=C_GREEN)
    txt(sl, h, cx2 + Inches(0.04), ty + Inches(0.06), cw, rh,
        size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
for i, row in enumerate(perf):
    name, rmse, da, dap, sharpe, col, note = row
    ry = ty + rh * (i + 1) + Inches(0.02)
    is_ens = "앙상블" in name
    fill = RGBColor(0x0A, 0x2A, 0x18) if is_ens else (C_PANEL if i%2==0 else C_DARK2)
    cells = [name, rmse, da, dap, sharpe, note]
    for j, (cell, cx2, cw) in enumerate(zip(cells, cxs, cws)):
        box(sl, cx2, ry, cw, rh - Inches(0.03), fill=fill)
        c = col if j == 0 else (C_GOLD if cell == "★" else C_WHITE)
        txt(sl, cell, cx2 + Inches(0.04), ry + Inches(0.06), cw - Inches(0.06), rh,
            size=12, bold=(j==0 and is_ens), color=c, align=PP_ALIGN.CENTER)

# 우측: 이항검정 설명
rx = Inches(7.35); rw = Inches(5.7)
box(sl, rx, Inches(1.62), rw, Inches(5.6), fill=C_PANEL, border=C_CYAN, bw=Pt(1.5))
txt(sl, "📌  이항검정 (Binomial Test)",
    rx + Inches(0.18), Inches(1.72), rw, Inches(0.38),
    size=14, bold=True, color=C_CYAN)

sections = [
    ("무엇인가?", C_GREEN, [
        "H₀: DA = 50% (동전던지기와 동일)",
        "H₁: DA > 50% (단측 검정)",
        "scipy.stats.binomtest 적용",
    ]),
    ("왜 필요한가?", C_GOLD, [
        "51~57% DA는 표본오차 범위일 수 있음",
        "p < 0.05 이어야 '유의미한 우위' 확인",
        "수치 신뢰도 정량 근거 확보",
    ]),
    ("결과 & 한계", C_RED, [
        "재학습 후 DA_p 자동 생성",
        "Walk-Forward 미적용 — 성능 과대평가 가능",
        "기말까지 WFV 적용 예정",
    ]),
]
sy = Inches(2.15)
for title, col, lines in sections:
    txt(sl, title, rx + Inches(0.18), sy, rw, Inches(0.32),
        size=13, bold=True, color=col)
    sy += Inches(0.35)
    for line in lines:
        txt(sl, f"  {line}", rx + Inches(0.18), sy, rw - Inches(0.3), Inches(0.3),
            size=12, color=C_LIGHT)
        sy += Inches(0.33)
    sy += Inches(0.12)

txt(sl,
    "p < 0.05 (★): 동전던지기 대비 통계적으로 유의한 방향 정확도",
    Inches(0.45), Inches(7.14), Inches(6.7), Inches(0.28),
    size=11, color=C_GREEN, italic=True)

# ── SL 6: 백테스트 시각화 ────────────────────────────────
sl = blank(prs); bg(sl)
section_label(sl, "B A C K T E S T   V I S U A L I Z A T I O N  ·  백 테 스 트   결 과")
slide_title(sl, "D+1 예측 vs 실제  |  최근 30 영업일 비교")
hline(sl, Inches(1.52))

_img = os.path.join(os.path.dirname(__file__), "outputs", "backtest_visual.png")
if os.path.exists(_img):
    # 이미지 비율 16:11 → 슬라이드 잔여 영역에 맞춤
    _iw, _ih = Inches(8.73), Inches(6.0)
    _il = (W - _iw) / 2
    sl.shapes.add_picture(_img, _il, Inches(1.45), _iw, _ih)
else:
    txt(sl, "⚠  outputs/backtest_visual.png 없음\n→  python visualize_backtest.py 를 먼저 실행하세요",
        Inches(2.5), Inches(3.3), Inches(8.3), Inches(0.8),
        size=14, color=C_RED, align=PP_ALIGN.CENTER)

txt(sl, "※ 검증셋 통계 기반 시뮬레이션  |  실운영 시 collect_data() 로 실제 시장 데이터 사용",
    Inches(0.45), Inches(7.2), W - Inches(0.9), Inches(0.22),
    size=9, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)

# ── SL 7: 핵심 인사이트 ──────────────────────────────────
sl = blank(prs); bg(sl)
section_label(sl, "K E Y   I N S I G H T S  ·  핵 심   인 사 이 트")
slide_title(sl, "핵심 인사이트")
hline(sl, Inches(1.52))

insights = [
    ("⚡", "앙상블이 RMSE 최소",
     "Ridge 앙상블 7.45원 — 단일 LGB(7.85원) 대비 5% 개선\nLSTM의 Sharpe 기여를 흡수해 최적 결합",
     C_GREEN, C_GREEN),
    ("📊", "LGB DA 57%가 최고",
     "단일 모델 중 방향 정확도 최고 (57.3%)\n단, 이항검정 p-value로 유의성 별도 검증 필요",
     C_CYAN, C_CYAN),
    ("🔗", "멀티스케일 피처 효과",
     "분봉(1H·5M·1M) 통계를 일봉에 통합한 137개 피처\n장중 변동성 정보가 일봉 예측에 실질 기여",
     C_GOLD, C_GOLD),
    ("🛡", "Persistence 비교 필수",
     "Persistence(기준선) RMSE 8.83원 → 앙상블 7.45원\n기준선 대비 15.6% 개선을 정량적으로 제시",
     RGBColor(0xA7, 0x8B, 0xFA), RGBColor(0xA7, 0x8B, 0xFA)),
]

iw = Inches(5.95); ih = Inches(2.1)
positions = [
    (Inches(0.45), Inches(1.62)),
    (Inches(6.55), Inches(1.62)),
    (Inches(0.45), Inches(3.85)),
    (Inches(6.55), Inches(3.85)),
]
for (lx, ty2), (emoji, title, desc, border, title_color) in zip(positions, insights):
    box(sl, lx, ty2, iw, ih, fill=C_PANEL, border=border, bw=Pt(1.5))
    txt(sl, f"{emoji}  {title}", lx + Inches(0.18), ty2 + Inches(0.12),
        iw - Inches(0.3), Inches(0.42), size=16, bold=True, color=title_color)
    txt(sl, desc, lx + Inches(0.18), ty2 + Inches(0.6),
        iw - Inches(0.3), Inches(1.3), size=13, color=C_LIGHT)

# 하단 wide 박스
box(sl, Inches(0.45), Inches(6.1), W - Inches(0.9), Inches(0.88),
    fill=C_DARK2, border=C_GREEN, bw=Pt(2))
txt(sl,
    "핵심 결론: 멀티스케일 피처 + 4모델 앙상블로 Persistence baseline 대비 15.6% RMSE 개선 "
    "— 단, Walk-Forward Validation으로 일반화 성능 추가 검증 필요",
    Inches(0.65), Inches(6.2), W - Inches(1.3), Inches(0.65),
    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── SL 8: 배포 & 대시보드 ─────────────────────────────────
sl = blank(prs); bg(sl)
section_label(sl, "D E P L O Y M E N T  ·  배 포   아 키 텍 처")
slide_title(sl, "실시간 배포 아키텍처")
hline(sl, Inches(1.52))

# 파이프라인 5단계
pipeline = [
    ("Google\nColab",    "T4 GPU\ntrain.py\n주 1회", C_GOLD),
    ("GitHub\n(main)",   "모델 파일\n코드 관리\nActions", C_GREEN),
    ("AWS S3",           "forecast.json\nresidual.csv\n백업", C_CYAN),
    ("EC2/Lambda",       "KST 06:00\npredict.py\n잔차수집", RGBColor(0xA7,0x8B,0xFA)),
    ("Streamlit\nCloud", "실시간 차트\nLGB 추론\n5분 캐시", C_RED),
]
pw2 = Inches(2.25); ph = Inches(2.3); py2 = Inches(1.68); px2 = Inches(0.45)
for i, (name, desc, col) in enumerate(pipeline):
    box(sl, px2, py2, pw2, ph, fill=C_PANEL, border=col, bw=Pt(2))
    txt(sl, name, px2, py2 + Inches(0.1), pw2, Inches(0.6),
        size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, desc, px2 + Inches(0.1), py2 + Inches(0.72), pw2 - Inches(0.18), Inches(1.4),
        size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, "→", px2 + pw2, py2 + Inches(0.85), Inches(0.22), Inches(0.5),
            size=20, bold=True, color=C_MUTED)
    px2 += pw2 + Inches(0.23)

# Colab 노트북 셀 구조
box(sl, Inches(0.45), Inches(4.18), W - Inches(0.9), Inches(3.0),
    fill=C_PANEL, border=C_MUTED, bw=Pt(1))
txt(sl, "colab_master.ipynb — 8셀 자동화",
    Inches(0.65), Inches(4.28), Inches(8), Inches(0.38),
    size=14, bold=True, color=C_GREEN)

cells = [
    ("셀 1", "GitHub Secrets", C_GOLD),
    ("셀 2", "Drive + 클론",   C_LIGHT),
    ("셀 3", "패키지 설치",     C_LIGHT),
    ("셀 4", "git pull",       C_LIGHT),
    ("셀 5", "train/predict", C_GREEN),
    ("셀 6", "로컬 테스트",    C_MUTED),
    ("셀 7", "git push",      C_CYAN),
    ("셀 8", "KST 06:00 자동",C_GOLD),
]
cw2 = Inches(1.5); ch = Inches(0.75); cy2 = Inches(4.72); cxs2 = []
cx2 = Inches(0.55)
for _ in cells[:4]:
    cxs2.append(cx2); cx2 += cw2 + Inches(0.12)
row2 = Inches(5.55)
cx2 = Inches(0.55); cxs2r = []
for _ in cells[4:]:
    cxs2r.append(cx2); cx2 += cw2 + Inches(0.12)

for i, ((no, label, col), lx) in enumerate(zip(cells[:4], cxs2)):
    box(sl, lx, cy2, cw2, ch, fill=C_DARK2, border=col, bw=Pt(1))
    txt(sl, no, lx + Inches(0.08), cy2 + Inches(0.06), Inches(0.5), Inches(0.3),
        size=12, bold=True, color=col)
    txt(sl, label, lx + Inches(0.08), cy2 + Inches(0.36), cw2 - Inches(0.14), Inches(0.3),
        size=11, color=C_WHITE)
for i, ((no, label, col), lx) in enumerate(zip(cells[4:], cxs2r)):
    box(sl, lx, row2, cw2, ch, fill=C_DARK2, border=col, bw=Pt(1))
    txt(sl, no, lx + Inches(0.08), row2 + Inches(0.06), Inches(0.5), Inches(0.3),
        size=12, bold=True, color=col)
    txt(sl, label, lx + Inches(0.08), row2 + Inches(0.36), cw2 - Inches(0.14), Inches(0.3),
        size=11, color=C_WHITE)

# 대시보드 피처 (우측)
dx = Inches(7.0)
box(sl, dx, Inches(4.18), Inches(6.1), Inches(3.0),
    fill=C_PANEL, border=C_GREEN, bw=Pt(1.5))
txt(sl, "Streamlit 대시보드 기능",
    dx + Inches(0.2), Inches(4.28), Inches(5.8), Inches(0.38),
    size=14, bold=True, color=C_GREEN)
dash_features = [
    ("TAB 1", "실시간 환율 차트 + 5분 자동 갱신",     C_GREEN),
    ("TAB 2", "D+1~D+22 예측 (LGB 실시간 추론)",      C_CYAN),
    ("TAB 3", "성능 비교 + Baseline + 유의성 표시",    C_GOLD),
    ("TAB 4", "거시지표 (VIX·DXY·WTI·금리차)",       C_LIGHT),
]
dfy = Inches(4.75)
for tab, desc, col in dash_features:
    txt(sl, f"▸  {tab}", dx + Inches(0.2), dfy, Inches(0.8), Inches(0.32),
        size=12, bold=True, color=col)
    txt(sl, desc, dx + Inches(1.05), dfy, Inches(4.9), Inches(0.32),
        size=12, color=C_LIGHT)
    dfy += Inches(0.42)

# ── SL 9: 결론 & 향후 계획 ───────────────────────────────
sl = blank(prs); bg(sl)
section_label(sl, "C O N C L U S I O N  ·  결 론")
slide_title(sl, "결론 & 향후 계획")
hline(sl, Inches(1.52))

# 왼쪽: 달성
box(sl, Inches(0.45), Inches(1.62), Inches(5.95), Inches(5.6),
    fill=C_PANEL, border=C_GREEN, bw=Pt(1.5))
txt(sl, "✅  달성한 것", Inches(0.65), Inches(1.72), Inches(5.7), Inches(0.38),
    size=15, bold=True, color=C_GREEN)

achieved = [
    ("RMSE 7.45원",         "Persistence(8.83원) 대비 15.6% 개선",     C_GREEN),
    ("4모델 앙상블",         "ARIMAX+LGB+LSTM+BiGRU → Ridge 최적 결합", C_GREEN),
    ("137개 멀티스케일 피처", "분봉 통계를 일봉에 통합 (9개 그룹)",       C_GREEN),
    ("이항검정 도입",         "DA_p — 교수 피드백 반영 유의성 검증",      C_CYAN),
    ("Persistence Baseline", "기준선 대비 정량 비교 체계 확립",           C_CYAN),
    ("End-to-End 자동화",    "Colab → GitHub → AWS → Streamlit",        C_GOLD),
    ("Lookahead Bias 차단",  "시간 역순 분할 + Train 전용 스케일러",     C_LIGHT),
]
ay = Inches(2.18)
for val, desc, col in achieved:
    txt(sl, f"▸  {val}", Inches(0.65), ay, Inches(2.3), Inches(0.32),
        size=13, bold=True, color=col)
    txt(sl, desc, Inches(3.0), ay, Inches(3.25), Inches(0.32),
        size=12, color=C_LIGHT)
    ay += Inches(0.42)

# 오른쪽: 향후 계획
box(sl, Inches(6.65), Inches(1.62), Inches(6.3), Inches(5.6),
    fill=C_PANEL, border=C_GOLD, bw=Pt(1.5))
txt(sl, "🎯  향후 계획 (기말까지)", Inches(6.85), Inches(1.72), Inches(6.0), Inches(0.38),
    size=15, bold=True, color=C_GOLD)

plans = [
    ("Walk-Forward Validation",
     "Expanding window 기반 이동 검증\nDA_p가 WFV에서도 유의한지 확인",
     C_GREEN),
    ("VIF 다중공선성 진단",
     "137개 피처 |r|>0.8 페어 분석\n노이즈 피처 정제로 모델 성능 개선",
     C_CYAN),
    ("레짐별 앙상블 (H3 가설)",
     "고변동/저변동 레짐 분리 후 별도 모델\n레짐별 RMSE 분해 정량 보고",
     C_GOLD),
    ("잔차 보정 모델",
     "residual_data.csv 30일 누적 후 학습\n24h 예측 오차의 반복 패턴 포착",
     RGBColor(0xA7, 0x8B, 0xFA)),
]
py3 = Inches(2.18)
for i, (title, desc, col) in enumerate(plans):
    box(sl, Inches(6.75), py3, Inches(6.1), Inches(1.22),
        fill=C_DARK2, border=col, bw=Pt(1.5))
    txt(sl, f"{i+1}.  {title}", Inches(6.95), py3 + Inches(0.1),
        Inches(5.8), Inches(0.38), size=14, bold=True, color=col)
    txt(sl, desc, Inches(6.95), py3 + Inches(0.5),
        Inches(5.8), Inches(0.6), size=12, color=C_LIGHT)
    py3 += Inches(1.35)

# 하단 교수 코멘트
box(sl, Inches(0.45), Inches(7.1), W - Inches(0.9), Inches(0.28), fill=C_BG)
txt(sl,
    '"이미 구축한 파이프라인의 결과가 얼마나 신뢰할 만한가를 정량적으로 검증하는 데 무게를 두어야 한다."  — 구자환 교수',
    Inches(0.5), Inches(7.12), W - Inches(1.0), Inches(0.28),
    size=11, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
out = "/home/user/usdkrw-prediction/usdkrw_presentation_v2.pptx"
prs.save(out)
print(f"✅  저장 완료: {out}  ({len(prs.slides)} 슬라이드)")
